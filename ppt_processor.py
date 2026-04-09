import os
import shutil
import subprocess
import tempfile
from typing import Dict
from pptx import Presentation


class PPTProcessor:
    def __init__(self):
        self.presentation = None
        self.output_path = None

    def load_presentation(self, ppt_path: str) -> bool:
        try:
            self.presentation = Presentation(ppt_path)
            print(f"Loaded presentation with {len(self.presentation.slides)} slides")
            return True
        except Exception as e:
            print(f"Error loading presentation: {e}")
            return False

    def keep_last_n_slides(self, n: int = 3):
        if not self.presentation:
            return False
        total_slides = len(self.presentation.slides)
        if total_slides <= n:
            print(f"Presentation has {total_slides} slides, keeping all.")
            return True
        slides_to_remove = total_slides - n
        print(f"Removing {slides_to_remove} slides, keeping last {n}")
        for _ in range(slides_to_remove):
            rId = self.presentation.slides._sldIdLst[0].rId
            self.presentation.part.drop_rel(rId)
            del self.presentation.slides._sldIdLst[0]
        return True

    def replace_response_boxes(self, mentee_response: str, parent_response: str):
        """
        Replace content inside the Mentee and Parent response boxes on the last slide.
        Identifies boxes by their placeholder text 'text', then sorts by horizontal
        position: leftmost = Mentee, rightmost = Parent.
        """
        if not self.presentation:
            return False

        last_slide = self.presentation.slides[-1]

        response_boxes = [
            shape for shape in last_slide.shapes
            if hasattr(shape, 'text_frame')
            and shape.text_frame.text.strip().lower() == 'text'
        ]

        if len(response_boxes) < 2:
            print(f"Warning: Expected 2 response boxes, found {len(response_boxes)}.")
            for shape in response_boxes:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.text.strip().lower() == 'text':
                            run.text = mentee_response
            return True

        response_boxes.sort(key=lambda s: s.left)
        mentee_box = response_boxes[0]
        parent_box = response_boxes[1]

        def set_box_text(shape, new_text):
            first = True
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if first:
                        run.text = new_text
                        first = False
                    else:
                        run.text = ''

        set_box_text(mentee_box, mentee_response)
        set_box_text(parent_box, parent_response)
        print(f"Set Mentee Response: '{mentee_response}'")
        print(f"Set Parent Response: '{parent_response}'")
        return True

    def save_as_pptx(self, output_path: str):
        if not self.presentation:
            return False
        try:
            self.presentation.save(output_path)
            self.output_path = output_path
            print(f"Saved presentation to: {output_path}")
            return True
        except Exception as e:
            print(f"Error saving presentation: {e}")
            return False

    def convert_to_pdf_linux(self, pptx_path: str, pdf_path: str) -> bool:
        """Convert PPTX to PDF using LibreOffice (works on Linux / cloud servers)."""
        try:
            out_dir = os.path.dirname(os.path.abspath(pdf_path))
            result = subprocess.run(
                ["libreoffice", "--headless", "--convert-to", "pdf",
                 "--outdir", out_dir, os.path.abspath(pptx_path)],
                capture_output=True, text=True, timeout=60
            )
            base = os.path.splitext(os.path.basename(pptx_path))[0]
            generated = os.path.join(out_dir, base + ".pdf")
            if os.path.exists(generated):
                shutil.move(generated, pdf_path)
                print(f"PDF saved to: {pdf_path}")
                return True
            print(f"LibreOffice output not found. stderr: {result.stderr}")
            return False
        except FileNotFoundError:
            print("LibreOffice not found. Falling back to PPTX copy.")
            alt_path = pdf_path.replace('.pdf', '_converted.pptx')
            shutil.copy2(pptx_path, alt_path)
            return True
        except Exception as e:
            print(f"PDF conversion error: {e}")
            return False

    def convert_to_pdf_windows(self, pptx_path: str, pdf_path: str) -> bool:
        """Convert PPTX to PDF on Windows using PowerPoint COM."""
        try:
            import win32com.client
            print("Converting to PDF using PowerPoint...")
            powerpoint = win32com.client.Dispatch("PowerPoint.Application")
            powerpoint.Visible = 1
            deck = powerpoint.Presentations.Open(os.path.abspath(pptx_path))
            deck.SaveAs(os.path.abspath(pdf_path), 32)
            deck.Close()
            powerpoint.Quit()
            print(f"PDF saved to: {pdf_path}")
            return True
        except Exception as e:
            print(f"Could not convert to PDF: {e}")
            alt_path = pdf_path.replace('.pdf', '_converted.pptx')
            shutil.copy2(pptx_path, alt_path)
            print(f"Saved as PPTX: {alt_path}")
            return True


def process_ppt_to_pdf(
    input_ppt_path: str,
    output_pdf_path: str,
    parent_response: str,
    mentee_response: str,
    keep_slides: int = 3
) -> bool:
    """Main function to process PPT file."""
    print("\n=== Starting PPT Processing ===")
    processor = PPTProcessor()

    if not processor.load_presentation(input_ppt_path):
        return False

    if not processor.keep_last_n_slides(keep_slides):
        return False

    if not processor.replace_response_boxes(mentee_response, parent_response):
        return False

    temp_dir = tempfile.gettempdir()
    temp_pptx_path = os.path.join(temp_dir, "temp_presentation.pptx")

    if not processor.save_as_pptx(temp_pptx_path):
        return False

    # Auto-detect platform: LibreOffice on Linux/cloud, PowerPoint on Windows
    import platform
    if platform.system() == "Windows":
        success = processor.convert_to_pdf_windows(temp_pptx_path, output_pdf_path)
    else:
        success = processor.convert_to_pdf_linux(temp_pptx_path, output_pdf_path)

    try:
        os.unlink(temp_pptx_path)
    except Exception:
        pass

    print("=== Processing Complete ===\n")
    return success
